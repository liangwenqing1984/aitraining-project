"""
生成"招聘数据智能采集与分析系统"介绍 PPT
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

OUTPUT_PATH = os.path.join(os.path.dirname(os.path.dirname(__file__)), "docs", "系统介绍.pptx")

# ========== 配色方案 ==========
PRIMARY = RGBColor(0x40, 0x9E, 0xFF)    # 主色蓝
DARK = RGBColor(0x30, 0x31, 0x33)        # 深色文字
GRAY = RGBColor(0x90, 0x93, 0x99)        # 灰色文字
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BG_DARK = RGBColor(0x1A, 0x1A, 0x2E)     # 深色背景
BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)    # 浅色背景
ACCENT_GREEN = RGBColor(0x67, 0xC2, 0x3A)
ACCENT_ORANGE = RGBColor(0xE6, 0xA2, 0x3C)
ACCENT_RED = RGBColor(0xF5, 0x6C, 0x6C)
ACCENT_PURPLE = RGBColor(0x9B, 0x59, 0xB6)
ACCENT_TEAL = RGBColor(0x00, 0xB8, 0xBA)
SLIDE_W = Inches(13.333)  # 16:9
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# ========== 辅助函数 ==========
def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_textbox(slide, left, top, width, height, text, font_size=18, color=DARK, bold=False, alignment=PP_ALIGN.LEFT, font_name='Microsoft YaHei'):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf

def add_multi_text(slide, left, top, width, height, lines, font_size=14, color=DARK, line_spacing=1.5):
    """lines: list of (text, bold, font_size_override, color_override)"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_info in enumerate(lines):
        if isinstance(line_info, str):
            text, bold, fs, clr = line_info, False, font_size, color
        else:
            text = line_info[0]
            bold = line_info[1] if len(line_info) > 1 else False
            fs = line_info[2] if len(line_info) > 2 else font_size
            clr = line_info[3] if len(line_info) > 3 else color
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(fs)
        p.font.color.rgb = clr
        p.font.bold = bold
        p.font.name = 'Microsoft YaHei'
        p.space_after = Pt(4)
    return tf

def add_card(slide, left, top, width, height, title, desc, title_color=PRIMARY, bg_color=None):
    """添加卡片式内容块"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color or RGBColor(0xFF, 0xFF, 0xFF)
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(12)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = title_color
    p.font.name = 'Microsoft YaHei'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.font.name = 'Microsoft YaHei'
    p2.space_before = Pt(4)

def add_accent_bar(slide, left, top, width, height, color=PRIMARY):
    """添加装饰色条"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

# ============================================================
# Slide 1: 封面
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BG_DARK)

# 顶部装饰线
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)

add_textbox(slide, 1.5, 1.8, 10.3, 1.0,
    '招聘数据智能采集与分析系统', font_size=40, color=WHITE, bold=True,
    alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 2.9, 10.3, 0.6,
    'AI-Powered Job Market Data Platform', font_size=20, color=PRIMARY,
    alignment=PP_ALIGN.CENTER)

add_textbox(slide, 1.5, 3.8, 10.3, 0.8,
    '一站式招聘数据采集、AI 增强、智能分析与语义检索平台',
    font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# 底部信息
add_accent_bar(slide, 0, 7.1, 13.333, 0.04, GRAY)
add_textbox(slide, 1.5, 7.2, 10.3, 0.3,
    '版本 2.0.0  |  2026.05', font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# Slide 2: 目录
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)
add_textbox(slide, 0.8, 0.4, 6, 0.6, '目录', font_size=32, color=DARK, bold=True)

toc_items = [
    ('01', '项目背景与痛点', '为什么需要这个系统'),
    ('02', '系统架构总览', '前后端 + 数据库 + 外部服务'),
    ('03', '数据采集', '智联/51job 双平台爬虫'),
    ('04', '反爬对抗体系', 'WAF 检测 + IP 代理池 + 指纹伪装'),
    ('05', 'AI 数据增强', 'LLM 标准化 8 个维度'),
    ('06', '智能分析', '7 维度图表 + AI 洞察报告'),
    ('07', '自然语言查询', 'Text-to-SQL 全流程'),
    ('08', 'RAG 语义搜索', 'pgvector + Ollama 向量化'),
    ('09', 'LLM 任务路由', '多模型智能调度'),
    ('10', '技术栈与总结', '前后端技术 + 未来展望'),
]

for i, (num, title, desc) in enumerate(toc_items):
    col = i % 2
    row = i // 2
    x = 0.8 + col * 6.2
    y = 1.4 + row * 1.1
    # 序号圆圈
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y + 0.05), Inches(0.4), Inches(0.4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PRIMARY
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.name = 'Microsoft YaHei'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    # 标题和描述
    add_textbox(slide, x + 0.6, y, 5, 0.3, title, font_size=16, color=DARK, bold=True)
    add_textbox(slide, x + 0.6, y + 0.3, 5, 0.25, desc, font_size=11, color=GRAY)

# ============================================================
# Slide 3: 项目背景与痛点
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)
add_textbox(slide, 0.8, 0.4, 10, 0.6, '01  项目背景与痛点', font_size=28, color=DARK, bold=True)

# 痛点卡片
pains = [
    ('🔍 数据获取难', '招聘网站反爬严格\nWAF/验证码/IP封锁层层设防\n手动采集效率低下'),
    ('📊 数据标准化难', '薪资格式千奇百怪\n"15K-20K·13薪"、"年薪30万"\n无法直接用于分析'),
    ('🧠 分析深度不够', '只看薪资分布不够\n需要技能趋势、行业对比\n多维度交叉分析'),
    ('🔎 精准搜索缺失', '关键词搜索太死板\n"后端开发"搜不到"Java工程师"\n语义相近但字面不同'),
]

for i, (title, desc) in enumerate(pains):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6 + i * 3.15), Inches(1.4), Inches(2.9), Inches(2.6)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFA, 0xFA, 0xFA)
    shape.line.color.rgb = RGBColor(0xE4, 0xE7, 0xED)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(14)
    tf.margin_right = Pt(14)
    tf.margin_top = Pt(14)
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = ACCENT_RED
    p.font.name = 'Microsoft YaHei'
    p2 = tf.add_paragraph()
    p2.text = desc
    p2.font.size = Pt(13)
    p2.font.color.rgb = DARK
    p2.font.name = 'Microsoft YaHei'
    p2.space_before = Pt(12)

# 箭头
add_textbox(slide, 0.8, 4.3, 11, 0.6,
    '⬇  解决方案：构建一站式 AI 驱动的招聘数据平台，从采集到分析全链路自动化',
    font_size=18, color=PRIMARY, bold=True)

# 解决亮点
solves = ['浏览器自动化\n爬虫 + IP代理池', 'LLM 标准化\n8 个维度增强', 'AI 生成\n专业洞察报告', 'pgvector\n语义向量搜索']
for i, s in enumerate(solves):
    add_card(slide, 0.6 + i * 3.15, 5.0, 2.9, 1.8, s.split('\n')[0], s.split('\n')[1],
             title_color=ACCENT_GREEN, bg_color=RGBColor(0xF0, 0xF9, 0xEB))

# ============================================================
# Slide 4: 系统架构总览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)
add_textbox(slide, 0.8, 0.4, 10, 0.6, '02  系统架构总览', font_size=28, color=DARK, bold=True)

# 架构图 - 用卡片表示各层
layers = [
    ('前端层 (Vue 3 + Element Plus)', '#409eff',
     ['数据采集 | 文件管理 | 智能分析 | 智能查询 | 语义搜索 | AI配置 | 系统管理']),
    ('API 网关 (Express + Socket.IO)', '#67c23a',
     ['REST API (60+ 端点) | WebSocket 实时推送 | JWT + OAuth2 认证 | RBAC 权限']),
    ('业务服务层 (16 个 Service)', '#e6a23c',
     ['爬虫服务 | LLM 服务 | 分析服务 | RAG 服务 | 代理池服务 | 用户/角色/权限服务']),
    ('数据层 (PostgreSQL + pgvector)', '#f56c6c',
     ['sp_tasks | sp_jobs | sp_job_enrichments | sp_job_embeddings | sp_llm_config | RBAC 表']),
    ('外部服务', '#9b59b6',
     ['Puppeteer (浏览器) | Ollama (本地 LLM) | jhao104/proxy_pool (代理池) | OAuth2 SSO']),
]

for i, (name, color_hex, descs) in enumerate(layers):
    y = 1.3 + i * 1.15
    color = RGBColor(int(color_hex[1:3], 16), int(color_hex[3:5], 16), int(color_hex[5:7], 16))
    add_accent_bar(slide, 0.8, y + 0.4, 0.08, 0.6, color)
    add_textbox(slide, 1.1, y, 3, 0.35, name, font_size=15, color=color, bold=True)
    add_textbox(slide, 1.1, y + 0.38, 11, 0.35, '  |  '.join(descs), font_size=12, color=GRAY)

# 右侧标注
add_textbox(slide, 9.5, 1.3, 3.5, 0.3, '⬅ Vite 代理转发', font_size=11, color=GRAY)
add_textbox(slide, 9.5, 2.45, 3.5, 0.3, '⬅ HTTP + WebSocket', font_size=11, color=GRAY)
add_textbox(slide, 9.5, 3.6, 3.5, 0.3, '⬅ better-sqlite3', font_size=11, color=GRAY)
add_textbox(slide, 9.5, 4.75, 3.5, 0.3, '⬅ pg + pgvector', font_size=11, color=GRAY)
add_textbox(slide, 9.5, 5.9, 3.5, 0.3, '⬅ HTTP / SDK', font_size=11, color=GRAY)

# ============================================================
# Slide 5: 核心功能概览
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)
add_textbox(slide, 0.8, 0.4, 10, 0.6, '核心功能总览 — 6 大模块', font_size=28, color=DARK, bold=True)

features = [
    ('数据采集', '智联/51job 双源\nIP 代理池自动切换\nWAF 智能绕过\n断点续传 + 重启恢复', ACCENT_GREEN),
    ('AI 数据增强', '薪资标准化 8 维度\n技能/分类/行业提取\nLLM 批量处理\n幂等 UPSERT', PRIMARY),
    ('智能分析', '7 维度可视化图表\nAI 生成洞察报告\nECharts 交互渲染\n历史报告管理', ACCENT_ORANGE),
    ('自然语言查询', '中文问题 → SQL\n安全白名单校验\nLLM 结果总结\n查询历史回顾', ACCENT_RED),
    ('RAG 语义搜索', '768 维向量索引\n余弦相似度匹配\n查询自动扩展\npgvector IVFFlat', ACCENT_PURPLE),
    ('IP 代理池', '11 个免费代理源\ncheckHealth 实时验证\n死代理自动淘汰\n支持扩展付费源', ACCENT_TEAL),
]

for i, (title, desc, color) in enumerate(features):
    col = i % 3
    row = i // 3
    x = 0.6 + col * 4.15
    y = 1.3 + row * 3.05
    # 卡片背景
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(3.85), Inches(2.75)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = WHITE
    shape.line.color.rgb = RGBColor(0xE4, 0xE7, 0xED)
    shape.line.width = Pt(1)
    # 顶部色条
    add_accent_bar(slide, x, y, 3.85, 0.06, color)
    # 标题
    add_textbox(slide, x + 0.25, y + 0.2, 3.35, 0.35, title, font_size=18, color=color, bold=True)
    # 描述
    add_textbox(slide, x + 0.25, y + 0.7, 3.35, 1.8, desc, font_size=13, color=DARK)

# ============================================================
# Slide 6: 数据采集
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)
add_textbox(slide, 0.8, 0.4, 10, 0.6, '03  数据采集 — 双平台爬虫', font_size=28, color=DARK, bold=True)

# 左侧：智联
add_accent_bar(slide, 0.6, 1.2, 5.9, 0.05, ACCENT_GREEN)
add_textbox(slide, 0.6, 1.35, 5.9, 0.35, '智联招聘 (zhilian.ts)', font_size=16, color=ACCENT_GREEN, bold=True)
zl_features = [
    '• 搜索 URL 查询参数编码 (?jl=622&kw=XX) 规避路径特征检测',
    '• 21 项 Chrome 启动参数 + 随机视口 + UA 轮换池 (4 种)',
    '• 独立 userDataDir 隔离会话 + 资源拦截保留 CSS',
    '• 详情页 axios API 直连优先 → 代理 → 浏览器渲染三级回退',
    '• WAF 空 body 检测 → 代理切换 → 降级串行 + 5-10s 大延迟',
    '• 互斥锁序列化 newPage + Promise.allSettled 最高并发 5',
    '• 每 5 组合主动重启浏览器刷新指纹 + 指数退避崩溃恢复',
]
add_multi_text(slide, 0.6, 1.8, 5.9, 3.8, zl_features, font_size=12)

# 右侧：51job
add_accent_bar(slide, 6.8, 1.2, 5.9, 0.05, ACCENT_ORANGE)
add_textbox(slide, 6.8, 1.35, 5.9, 0.35, '前程无忧 (job51.ts)', font_size=16, color=ACCENT_ORANGE, bold=True)
j5_features = [
    '• puppeteer-extra + StealthPlugin 专业反检测',
    '• 每页独立指纹：视口/UA/plugins/chrome/permissions 伪造',
    '• XHR 拦截主力提取 (we.51job.com/api/job/search-pc)',
    '• WAF 双重防线：硬编码 5 类签名 (零延迟) + AI 页面分类 (语义)',
    '• 三级 WAF 恢复：重载(10s) → 长等(45-90s) → 首页回退',
    '• 关键参数 reportType=1 (业务报表绕过 Aliyun WAF)',
    '• 首次访问前首页 Cookie 建立 (3-7s 停留 + 滚动)',
]
add_multi_text(slide, 6.8, 1.8, 5.9, 3.8, j5_features, font_size=12)

# 底部对比总结
add_accent_bar(slide, 0.6, 5.8, 12.1, 0.03, GRAY)
add_textbox(slide, 0.8, 5.95, 11.7, 0.8,
    '对比总结：智联核心在 API 直连逃逸 WAF，51job 核心在 XHR 拦截 + 双重 WAF 检测 + reportType=1 参数绕过。两者均支持断点续传 (JSONB) + 最大 10 次浏览器重启。',
    font_size=12, color=GRAY)

# ============================================================
# Slide 7: 反爬对抗体系
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)
add_textbox(slide, 0.8, 0.4, 10, 0.6, '04  反爬对抗体系', font_size=28, color=DARK, bold=True)

anti_crawl = [
    ('🖥️ 浏览器层伪装',
     ['随机视口 1366-1920×768-1080', 'UA 轮换池 4 种', '独立 userDataDir', '21 项 Chrome flags',
      '每页独立指纹 (plugins/chrome)', 'puppeteer-extra StealthPlugin']),
    ('🕐 类人行为模拟',
     ['分级延迟 2-10s', '渐进式懒加载滚动 (8次×800ms)', '批次间隔 8-10s',
      '首页 Cookie 建立 (停留+滚动)', '批次内错峰 0.8-2s 偏移']),
    ('🛡️ WAF 检测',
     ['AI 页面分类 6 种类型', '硬编码签名检测 (零延迟)', '智联：API 直连逃逸',
      '51job：reportType=1 绕过', '5s 冷却防高频 LLM 调用']),
    ('🔄 恢复策略',
     ['51job 三级恢复：重载→长等→回退', '智联：代理切换→降级串行', '指数退避 30s×2ⁿ (上限120s)',
      '浏览器崩溃断点续传', '最大 10 次重启保护']),
]

for i, (title, items) in enumerate(anti_crawl):
    x = 0.5 + i * 3.2
    add_accent_bar(slide, x, 1.2, 3.0, 0.05, [ACCENT_GREEN, ACCENT_ORANGE, ACCENT_RED, PRIMARY][i])
    add_textbox(slide, x, 1.35, 3.0, 0.35, title, font_size=14, color=DARK, bold=True)
    for j, item in enumerate(items):
        add_textbox(slide, x + 0.1, 1.8 + j * 0.35, 2.9, 0.3, f'• {item}', font_size=11, color=DARK)

# 底部架构
add_accent_bar(slide, 0.6, 4.3, 12.1, 0.03, GRAY)
add_textbox(slide, 0.8, 4.45, 11.7, 0.3, '协同架构', font_size=14, color=DARK, bold=True)
add_textbox(slide, 0.8, 4.8, 11.7, 1.5,
    'proxy_pool :5010  ─→  ProxyPool.getProxy()  ─→  checkHealth(zhaopin.com/51job.com)  ─→  浏览器 --proxy-server\n'
    '                                                                              ↓ 失败\n'
    '                                                                   deleteProxy() → 换新代理\n'
    '                                                                  池耗尽 → 直连降级',
    font_size=12, color=GRAY)

# ============================================================
# Slide 8: AI 数据增强
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)
add_textbox(slide, 0.8, 0.4, 10, 0.6, '05  AI 数据增强', font_size=28, color=DARK, bold=True)

# 8 个增强维度
dims = [
    ('💰 薪资标准化', '15K-20K·13薪\n→ 15000-20000/月'),
    ('🏷️ 职位分类', '14大类+子类\n技术→后端开发'),
    ('🏭 行业识别', '互联网/金融/制造\n等 14 类标准'),
    ('🔧 技能提取', '必备+加分分离\n技术栈关联'),
    ('🎓 学历规范', '5 级归一化\n本科/硕士/博士'),
    ('📅 经验年限', '3-5年→3-5\n应届→0-1'),
    ('🎁 福利识别', '五险一金/年终奖\n关键词识别'),
    ('🏠 工作模式', '远程/现场/混合\n三分类'),
]

for i, (title, desc) in enumerate(dims):
    col = i % 4
    row = i // 4
    add_card(slide, 0.5 + col * 3.15, 1.2 + row * 2.0, 2.9, 1.75, title, desc,
             title_color=PRIMARY)

# 技术特点
add_accent_bar(slide, 0.6, 5.4, 12.1, 0.03, GRAY)
add_textbox(slide, 0.8, 5.55, 11.7, 1.2,
    '技术特点：BATCH_SIZE=1 逐条处理 + 500ms 间隔避免 API 限流  |  3 次重试 + 递增 temperature 提高成功率\n'
    '3 层降级 JSON 解析 (直接解析 → 正则提取 → 单引号修复)  |  ON CONFLICT UPSERT 幂等可重跑\n'
    'WebSocket 实时推送增强进度  |  支持多模型路由 (enrichment 任务类型)',
    font_size=12, color=GRAY)

# ============================================================
# Slide 9: 智能分析
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)
add_textbox(slide, 0.8, 0.4, 10, 0.6, '06  智能分析与洞察报告', font_size=28, color=DARK, bold=True)

# 7 维度图表
add_textbox(slide, 0.8, 1.2, 5, 0.35, '7 维度数据可视化', font_size=16, color=DARK, bold=True)
charts = ['薪资分布 (箱线图/柱状图)', '城市分布 (饼图/地图)', '学历要求分布 (柱状图)',
          '经验要求分布 (柱状图)', '行业分布 (饼图)', '技能词云 (云图)', '工作模式 (饼图)']
for i, c in enumerate(charts):
    col = i % 4
    row = i // 4
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.7 + col * 2.2), Inches(1.7 + row * 0.55), Inches(2.0), Inches(0.42)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xFD, 0xF6, 0xEC)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.paragraphs[0].text = c
    tf.paragraphs[0].font.size = Pt(11)
    tf.paragraphs[0].font.color.rgb = ACCENT_ORANGE
    tf.paragraphs[0].font.name = 'Microsoft YaHei'
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# AI 洞察报告
add_textbox(slide, 0.8, 3.2, 5, 0.35, 'AI 洞察报告生成', font_size=16, color=DARK, bold=True)
add_multi_text(slide, 0.8, 3.6, 6, 2.5, [
    '1. 从 job_enrichments 聚合多维度统计数据',
    '2. 注入 LLM Prompt 模板，含完整统计表',
    '3. LLM 返回结构化 JSON：',
    '   { title, summary, sections[{heading, body, key_findings}] }',
    '4. charts_config 为 ECharts option 对象，前端直接渲染',
    '5. 失败降级：JSON 解析失败 → 正则提取 Markdown',
], font_size=12)

# 右侧架构简图
add_textbox(slide, 7.5, 1.2, 5.3, 5.0,
    '数据流\n\n'
    'sp_job_enrichments\n'
    '    ↓ SQL 聚合\n'
    '薪资/城市/技能/行业\n'
    '学历/经验/模式 统计\n'
    '    ↓ LLM Prompt\n'
    '结构化 JSON\n'
    '    ↓ 解析 + 入库\n'
    'sp_market_reports\n'
    '    ↓ WebSocket\n'
    '前端 ECharts 渲染\n'
    'Markdown 报告展示',
    font_size=12, color=GRAY)

# ============================================================
# Slide 10: 自然语言查询
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)
add_textbox(slide, 0.8, 0.4, 10, 0.6, '07  自然语言查询 (Text-to-SQL)', font_size=28, color=DARK, bold=True)

add_textbox(slide, 0.8, 1.2, 11, 2.2,
    '用户输入：  "北京 Java 岗位薪资 20K 以上的有哪些" \n'
    '                    ↓\n'
    '① 意图解析：LLM 提取实体 (Java/北京/20K) + 聚合目标\n'
    '② Schema 注入：System Prompt 含 job_enrichments 完整表结构 (字段/类型/枚举值)\n'
    '③ SQL 生成：LLM 输出 PostgreSQL 查询 (SELECT/JOIN/GROUP BY/子查询)\n'
    '④ 安全校验：白名单正则 — 仅允许 SELECT，拦截 DROP/INSERT；自动 LIMIT 500\n'
    '⑤ 执行查询：pool.query(sql) 参数化执行\n'
    '⑥ 结果总结：LLM 用 2-3 句中文总结关键数值和趋势',
    font_size=13, color=DARK)

add_accent_bar(slide, 0.6, 3.8, 12.1, 0.03, GRAY)

# 示例
add_textbox(slide, 0.8, 4.0, 6, 0.35, '快捷示例', font_size=16, color=DARK, bold=True)
examples = ['"薪资最高的 10 个岗位是哪些？"',
            '"各城市 Java 岗位平均薪资对比"',
            '"本科学历要求的岗位有多少个？"',
            '"互联网行业平均薪资范围"',
            '"哈尔滨和沈阳前端岗位数量和薪资对比"']
for i, ex in enumerate(examples):
    add_textbox(slide, 0.8, 4.5 + i * 0.45, 5.5, 0.35, f'• {ex}', font_size=12, color=DARK)

# 安全机制
add_textbox(slide, 7.5, 4.0, 5.3, 0.35, '安全白名单机制', font_size=16, color=ACCENT_RED, bold=True)
add_multi_text(slide, 7.5, 4.5, 5.3, 2.5, [
    '• 仅允许 SELECT/WITH 语句',
    '• 拦截 INSERT/DROP/TRUNCATE/DELETE',
    '• 自动追加 LIMIT 500',
    '• 多语句分号截断防护',
    '• saved_queries 持久化历史',
], font_size=12)

# ============================================================
# Slide 11: RAG 语义搜索
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)
add_textbox(slide, 0.8, 0.4, 10, 0.6, '08  RAG 语义搜索', font_size=28, color=DARK, bold=True)

add_textbox(slide, 0.8, 1.2, 11, 2.5,
    '完整流水线\n\n'
    'sp_job_enrichments (AI增强后结构化数据)\n'
    '        ↓  buildJobText() — 拼接职位名/公司/城市/技能/行业为自然语言段落\n'
    '  "职位名称：Java高级工程师；公司：X科技；城市：北京；技能：Java, Spring, MySQL"\n'
    '        ↓  Ollama nomic-embed-text → 768 维浮点向量 (POST /api/embeddings)\n'
    '        ↓  INSERT INTO sp_job_embeddings (embedding) VALUES ($1::vector)\n'
    '        ↓  IVFFlat 索引 (100 lists, cosine_ops) — 近似搜索比全量 KNN 快 10-100 倍\n'
    'semanticSearch() → 1 - (embedding <=> query) = 余弦相似度 (0-1) → ORDER BY similarity DESC',
    font_size=13, color=DARK)

add_accent_bar(slide, 0.6, 4.2, 12.1, 0.03, GRAY)

# 查询扩展
add_textbox(slide, 0.8, 4.4, 6, 0.35, '查询扩展 (Query Expansion)', font_size=16, color=DARK, bold=True)
add_multi_text(slide, 0.8, 4.8, 5.5, 2.0, [
    '• 短查询 (≤10字符) 自动触发',
    '• 30+ 术语映射表：Java → "Java开发 Spring 微服务"',
    '• 前端 → "Vue React Web前端 网页开发"',
    '• 解决短查询语义稀疏问题',
    '• 长查询 (>10字符) 不做扩展，保留原意',
], font_size=12)

# 右侧
add_textbox(slide, 7.5, 4.4, 5.3, 0.35, '关键参数', font_size=16, color=DARK, bold=True)
add_multi_text(slide, 7.5, 4.8, 5.3, 2.0, [
    '• 模型：nomic-embed-text (768维)',
    '• 存储：pgvector vector(768)',
    '• 索引：IVFFlat 100 lists',
    '• 相似度：默认阈值 0.3',
    '• 间隔：200ms 避免 Ollama 过载',
    '• 幂等：ON CONFLICT UPSERT',
], font_size=12)

# ============================================================
# Slide 12: LLM 任务路由
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)
add_textbox(slide, 0.8, 0.4, 10, 0.6, '09  LLM 任务路由', font_size=28, color=DARK, bold=True)

add_textbox(slide, 0.8, 1.2, 11.7, 0.5,
    '系统支持同时配置多个 AI 模型，不同任务类型自动选择对应模型。核心机制：每个 LLM 配置维护 task_routing JSONB 数组。',
    font_size=14, color=DARK)

# 四种任务类型
task_types = [
    ('enrichment\n数据增强', 'DeepSeek / Ollama', '逐条标准化职位数据\n批量处理，成本敏感'),
    ('insights\n智能洞察', 'DeepSeek / GPT-4o', '聚合统计 + 生成报告\n强推理，长文输出'),
    ('query\nNL 查询', 'DeepSeek / 智谱 GLM', '自然语言 → SQL\n中文 NLU 强'),
    ('anti-crawl\n反爬检测', 'Ollama qwen3:4b', '页面分类 + 选择器推荐\n低延迟 (<1s)'),
]

for i, (task, model, desc) in enumerate(task_types):
    x = 0.5 + i * 3.2
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(1.9), Inches(3.0), Inches(2.0)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xF5, 0xF7, 0xFA)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(10)
    tf.margin_top = Pt(8)
    p = tf.paragraphs[0]
    p.text = task
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = PRIMARY
    p.font.name = 'Microsoft YaHei'
    p2 = tf.add_paragraph()
    p2.text = f'推荐：{model}'
    p2.font.size = Pt(11)
    p2.font.color.rgb = GRAY
    p2.font.name = 'Microsoft YaHei'
    p2.space_before = Pt(6)
    p3 = tf.add_paragraph()
    p3.text = desc
    p3.font.size = Pt(11)
    p3.font.color.rgb = DARK
    p3.font.name = 'Microsoft YaHei'
    p3.space_before = Pt(4)

# 路由逻辑
add_textbox(slide, 0.8, 4.2, 11.7, 0.35, '路由选择逻辑', font_size=16, color=DARK, bold=True)
add_textbox(slide, 0.8, 4.6, 11.7, 2.2,
    'getConfigForTask(taskType)\n'
    '  ├─ refreshConfigCache() → 从 sp_llm_config 加载 active 配置 (60s 缓存)\n'
    '  ├─ 遍历 configs，检查 taskRouting JSONB 数组是否包含 taskType\n'
    '  ├─ 返回首个匹配 → 初始化 Provider (CloudProvider / LocalProvider)\n'
    '  └─ 无匹配 → 兜底返回第一个 active 配置 (保证始终有可用模型)\n\n'
    'API Key 加密：AES-256-GCM 加密存储，格式自动检测 (isEncrypted 正则)，明文 Key 首次使用自动升级',
    font_size=12, color=DARK)

# ============================================================
# Slide 13: IP 代理池
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)
add_textbox(slide, 0.8, 0.4, 10, 0.6, 'IP 代理池', font_size=28, color=DARK, bold=True)

add_textbox(slide, 0.8, 1.1, 11.7, 1.2,
    '外部 proxy_pool (jhao104) :5010         本项目 ProxyPool 类                   爬虫使用\n'
    '┌──────────────────────┐       ┌──────────────────────┐       ┌──────────────────┐\n'
    '│ GET /get/ → 随机代理  │  ───→ │ getProxy() + checkHealth │  ───→ │ 智联: --proxy-server│\n'
    '│ GET /delete/ → 删除   │  ←─── │ deleteProxy()         │       │ 51job: axios 回退  │\n'
    '│ GET /count → 代理总数  │       │ 连续 5 次失败停止返回   │       │ deadProxyCache    │\n'
    '└──────────────────────┘       └──────────────────────┘       └──────────────────┘',
    font_size=12, color=DARK)

add_accent_bar(slide, 0.6, 2.6, 12.1, 0.03, GRAY)

# 左：11 个免费源
add_textbox(slide, 0.8, 2.8, 6, 0.35, '11 个免费代理源', font_size=14, color=DARK, bold=True)
sources = ['站大爷 (zdaye.com) — HTML xpath 解析', '66代理 (66ip.cn) — xpath',
           '开心代理 (kxdaili.com) — xpath', 'FreeProxyList — URL解码+正则',
           '快代理 (kuaidaili.com) — xpath 多页', '冰凌代理 (binglx.cn) — 更新最快 ★★★',
           '云代理 (ip3366.net) — 正则', '小幻代理 (ip.ihuan.me) — 正则',
           '免费代理库 (jiangxianli.com)', '89免费代理 (89ip.cn) — 正则',
           '稻壳代理 (docip.net) — JSON API 可用率最高 ★★★']
for i, s in enumerate(sources):
    col = i % 2
    row = i // 2
    add_textbox(slide, 0.8 + col * 3.1, 3.25 + row * 0.28, 3.0, 0.25, f'• {s}', font_size=9, color=DARK)

# 右：关键机制
add_textbox(slide, 7.5, 2.8, 5.3, 0.35, '关键机制', font_size=14, color=DARK, bold=True)
add_multi_text(slide, 7.5, 3.25, 5.3, 3.5, [
    'checkHealth() 实时验证',
    '  • 8s 超时 + 不跟随重定向',
    '  • 仅 2xx 视为可用',
    '  • 最多重试 3 次',
    '',
    '死代理自动淘汰',
    '  • ERR_TUNNEL_CONNECTION_FAILED',
    '  • deleteProxy() 从池中删除',
    '  • deadProxyCache Set 黑名单',
    '',
    '降级直连',
    '  • 代理池耗尽 → 浏览器直连',
    '  • 切换次数达上限 → 直连',
    '',
    '扩展付费代理只需两步',
    '  1. proxyFetcher.py 添加 @staticmethod',
    '  2. setting.py PROXY_FETCHER 注册',
], font_size=11)

# ============================================================
# Slide 14: 技术栈
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)
add_textbox(slide, 0.8, 0.4, 10, 0.6, '10  技术栈', font_size=28, color=DARK, bold=True)

stacks = [
    ('前端', ACCENT_GREEN, ['Vue 3 (Composition API)', 'Element Plus', 'Pinia 状态管理',
                             'Vue Router', 'ECharts 6', 'Socket.IO Client', 'Axios', 'Vite', 'marked']),
    ('后端', PRIMARY, ['Node.js + Express', 'TypeScript', 'Puppeteer (浏览器自动化)',
                        'Socket.IO (WebSocket)', 'ExcelJS (Excel 读写)', 'better-sqlite3 / pg',
                        'AES-256-GCM 加密', 'bcryptjs 密码哈希']),
    ('数据', ACCENT_ORANGE, ['PostgreSQL (SeaboxSQL)', 'pgvector 向量扩展',
                              'Redis (代理池存储)', 'JSONB 灵活配置', 'IVFFlat 向量索引']),
    ('AI / LLM', ACCENT_PURPLE, ['DeepSeek v4-pro', 'OpenAI GPT-4o', 'Ollama (本地推理)',
                                  '智谱 GLM', 'Anthropic Claude', 'nomic-embed-text (768维)']),
    ('运维', ACCENT_TEAL, ['PM2 进程管理', 'Nginx 反向代理', 'jhao104/proxy_pool',
                            'OAuth2 统一认证', '定时备份 (pg_dump)']),
]

for i, (name, color, items) in enumerate(stacks):
    x = 0.5 + i * 2.55
    add_accent_bar(slide, x, 1.1, 2.3, 0.05, color)
    add_textbox(slide, x, 1.25, 2.3, 0.35, name, font_size=14, color=color, bold=True)
    for j, item in enumerate(items):
        add_textbox(slide, x + 0.1, 1.7 + j * 0.32, 2.3, 0.28, f'• {item}', font_size=10, color=DARK)

# 底部数据库
add_accent_bar(slide, 0.6, 4.8, 12.1, 0.03, GRAY)
add_textbox(slide, 0.8, 4.95, 11.7, 0.3, '数据库 15 张表：', font_size=13, color=DARK, bold=True)
tables = '业务表: sp_tasks | sp_csv_files | sp_jobs | sp_job_enrichments | sp_market_reports | sp_saved_queries | sp_job_embeddings | sp_llm_config'
tables2 = 'RBAC表: sp_users | sp_roles | sp_permissions | sp_menus | sp_user_roles | sp_role_permissions | sp_role_menus'
add_textbox(slide, 0.8, 5.3, 11.7, 0.3, tables, font_size=9, color=GRAY)
add_textbox(slide, 0.8, 5.55, 11.7, 0.3, tables2, font_size=9, color=GRAY)

# ============================================================
# Slide 15: 总结与展望
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_DARK)
add_accent_bar(slide, 0, 0, 13.333, 0.06, PRIMARY)

add_textbox(slide, 0.8, 0.6, 11.7, 0.6, '总结与展望', font_size=36, color=WHITE, bold=True)

# 核心数据
add_textbox(slide, 0.8, 1.5, 11.7, 0.35, '── 核心数据 ──', font_size=16, color=PRIMARY)
metrics = [('60+', 'API 端点'), ('16', '服务模块'), ('15', '数据库表'), ('11', '代理源'),
           ('8', '增强维度'), ('7', '分析图表'), ('4', '任务路由'), ('2', '爬虫平台')]
for i, (num, label) in enumerate(metrics):
    x = 0.8 + i * 1.55
    add_textbox(slide, x, 2.0, 1.4, 0.6, num, font_size=32, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, 2.55, 1.4, 0.3, label, font_size=11, color=GRAY, alignment=PP_ALIGN.CENTER)

add_accent_bar(slide, 0.8, 3.0, 11.7, 0.02, GRAY)

# 已实现
add_textbox(slide, 0.8, 3.2, 11.7, 0.35, '✅ 已实现', font_size=16, color=ACCENT_GREEN)
add_multi_text(slide, 0.8, 3.6, 5.5, 3.0, [
    '• 智联/51job 双平台自动化采集',
    '• AI 驱动 8 维度数据增强',
    '• 7 维度图表 + AI 洞察报告',
    '• 自然语言 → SQL 查询',
    '• pgvector 语义向量搜索',
    '• IP 代理池自动切换',
    '• WAF 双重检测 + 三级恢复',
    '• LLM 多模型智能路由',
    '• RBAC 用户/角色/权限管理',
    '• WebSocket 实时进度推送',
    '• 断点续传 + 浏览器崩溃恢复',
    '• AES-256-GCM API Key 加密',
], font_size=12, color=WHITE)

# 展望
add_textbox(slide, 7.5, 3.2, 5.3, 0.35, '🔮 未来展望', font_size=16, color=ACCENT_ORANGE)
add_multi_text(slide, 7.5, 3.6, 5.3, 3.0, [
    '• 扩展更多招聘平台 (Boss直聘等)',
    '• 接入付费代理源提升稳定性',
    '• 引入消息队列 (Redis/RabbitMQ)',
    '• 爬虫分布式部署 + 任务调度',
    '• 数据大屏实时可视化',
    '• 移动端适配 (PWA)',
    '• 自动定时采集 + 邮件报告',
    '• CI/CD 自动化部署',
    '• 接入更多 LLM 模型',
    '• 技能知识图谱构建',
], font_size=12, color=WHITE)

# 底部
add_accent_bar(slide, 0, 7.1, 13.333, 0.04, GRAY)
add_textbox(slide, 0.8, 7.2, 11.7, 0.3, '谢谢！', font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# ========== 保存 ==========
prs.save(OUTPUT_PATH)
print(f'PPT 已生成: {OUTPUT_PATH}')
